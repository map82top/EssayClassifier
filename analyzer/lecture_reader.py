from pptx import Presentation
import re
import validators


def read_from_file(path):
    prs = Presentation(path)
    text = ''
    for slide in prs.slides:
        text += handle_slide(slide)
    return text.strip()


def read_from_presentation(presentation):
    text = ''
    for slide in presentation.slides:
        text += handle_slide(slide)
    return text.strip()


def handle_slide(slide):
    text = ''
    for element in slide.shapes:
        if hasattr(element, "text"):
            for word in re.split(' ', element.text):
                if not validators.url(word) and word.strip() != '':
                    text += ' ' + word

    return text
